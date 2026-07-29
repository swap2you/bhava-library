"""Technical metadata extraction and coverage tests."""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter

from bhava_library.curation.enrich import (
    extract_technical_metadata,
    summarize_metadata_coverage,
)


def test_pdf_image_and_office_metadata(tmp_path: Path) -> None:
    pdf = tmp_path / "sample.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    writer.add_metadata({"/Title": "Test title"})
    with pdf.open("wb") as stream:
        writer.write(stream)
    pdf_result = extract_technical_metadata(pdf)
    assert pdf_result["extraction_status"] == "full"
    assert pdf_result["pdf"]["page_count"] == 1
    assert pdf_result["pdf"]["searchable_text_available"] is False
    assert pdf_result["pdf"]["likely_scanned"] is True

    image_path = tmp_path / "image.png"
    Image.new("RGB", (13, 17)).save(image_path)
    image_result = extract_technical_metadata(image_path)
    assert image_result["image"]["width"] == 13
    assert image_result["image"]["height"] == 17
    assert image_result["image"]["mode"] == "RGB"

    docx_path = tmp_path / "document.docx"
    document = Document()
    document.core_properties.title = "Safe document"
    document.save(docx_path)
    assert extract_technical_metadata(docx_path)["office"]["core_properties"]["title"] == (
        "Safe document"
    )

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    presentation.save(pptx_path)
    assert extract_technical_metadata(pptx_path)["office"]["slide_count"] == 1

    xlsx_path = tmp_path / "sheets.xlsx"
    workbook = Workbook()
    workbook.create_sheet("Second")
    workbook.save(xlsx_path)
    workbook.close()
    assert extract_technical_metadata(xlsx_path)["office"]["sheet_count"] == 2


def test_zip_metadata_only_and_quarantine_skip(tmp_path: Path) -> None:
    archive = tmp_path / "bundle.zip"
    with zipfile.ZipFile(archive, "w") as zf:
        zf.writestr("lesson.txt", "lesson")
        executable = zipfile.ZipInfo("run.sh")
        executable.external_attr = 0o100755 << 16
        zf.writestr(executable, "#!/bin/sh")

    result = extract_technical_metadata(archive)
    assert result["extraction_status"] == "full"
    assert result["archive"]["member_count"] == 2
    assert result["archive"]["uncompressed_bytes"] == 15
    assert result["archive"]["has_executable_members"] is True
    assert result["archive"]["has_encrypted_members"] is False

    quarantined = extract_technical_metadata(archive, quarantined=True)
    assert quarantined["extraction_status"] == "skipped_quarantine"
    assert "archive" not in quarantined

    quarantined_tar = tmp_path / "quarantined" / "bundle.tar"
    quarantined_tar.parent.mkdir()
    quarantined_tar.write_bytes(b"not opened")
    assert extract_technical_metadata(quarantined_tar)["extraction_status"] == "skipped_quarantine"


def test_fallback_errors_availability_and_coverage(tmp_path: Path, monkeypatch) -> None:
    unknown = tmp_path / "record.bin"
    unknown.write_bytes(b"metadata only")
    monkeypatch.setattr(
        "bhava_library.curation.enrich.extractor_availability",
        lambda: {
            "pypdf": False,
            "mutagen": False,
            "Pillow": False,
            "python-docx": False,
            "python-pptx": False,
            "openpyxl": False,
        },
    )
    fallback = extract_technical_metadata(unknown)
    assert fallback["extraction_status"] == "fallback_only"
    assert fallback["extractor_availability"]["pypdf"] is False

    missing_pdf = extract_technical_metadata(tmp_path / "missing.pdf")
    assert missing_pdf["extraction_status"] == "error"
    assert missing_pdf["errors"]

    report = summarize_metadata_coverage(
        [
            {"extraction_status": "full"},
            {"extraction_status": "partial"},
            fallback,
            missing_pdf,
        ]
    )
    assert json.loads(json.dumps(report))["technical_metadata_count"] == 2
    assert report["coverage_percent"] == 50.0
